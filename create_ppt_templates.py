#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建风格一致的PPT模板
基于vivo AIGC赛事PPT风格：
- 主色调：深紫蓝 #6B3FA0
- 背景：深色 #1A1A2E（封面）/ 白色 #FFFFFF（内容页）
- 字体：微软雅黑 + Arial
- 16:9 宽屏
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 定义颜色
PRIMARY_COLOR = RGBColor(0x6B, 0x3F, 0xA0)      # 深紫蓝
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)            # 深色背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)              # 白色
LIGHT_PURPLE = RGBColor(0x9B, 0x6F, 0xD0)       # 浅紫色
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xE2)        # 强调蓝
ACCENT_PINK = RGBColor(0xE2, 0x4A, 0x8A)        # 强调粉
ACCENT_GREEN = RGBColor(0x4A, 0xE2, 0x9B)       # 强调绿
ACCENT_ORANGE = RGBColor(0xE2, 0x8A, 0x4A)      # 强调橙

def set_slide_size(prs):
    """设置16:9宽屏"""
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle="", accent_color=PRIMARY_COLOR):
    """创建封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_color
    top_bar.line.fill.background()
    
    # 左侧装饰块
    left_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(1.5), Inches(0.3), Inches(2.5)
    )
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = accent_color
    left_block.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = "Microsoft YaHei"
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(12), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.name = "Microsoft YaHei"
    
    # 底部装饰线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.5), Inches(4), Inches(0.04)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = accent_color
    bottom_line.line.fill.background()
    
    return slide

def add_content_slide(prs, title, content_items, accent_color=PRIMARY_COLOR):
    """创建内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # 左侧紫色块
    left_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(3.5), prs.slide_height
    )
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = accent_color
    left_block.line.fill.background()
    
    # 左侧标题
    left_title = slide.shapes.add_textbox(Inches(0.4), Inches(2.2), Inches(2.7), Inches(1))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    
    # 左侧英文
    left_en = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(2.7), Inches(0.5))
    tf = left_en.text_frame
    p = tf.paragraphs[0]
    p.text = "CONTENTS"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    
    # 左侧装饰线
    left_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(3.7), Inches(1.5), Inches(0.03)
    )
    left_line.fill.solid()
    left_line.fill.fore_color.rgb = WHITE
    left_line.line.fill.background()
    
    # 右侧内容
    y_pos = Inches(1.5)
    for i, item in enumerate(content_items, 1):
        # 序号
        num_box = slide.shapes.add_textbox(Inches(4), y_pos, Inches(0.8), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"0{i}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.font.name = "Arial"
        
        # 内容标题
        content_box = slide.shapes.add_textbox(Inches(5), y_pos, Inches(7.5), Inches(0.6))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.font.name = "Microsoft YaHei"
        
        y_pos += Inches(1.2)
    
    return slide

def add_section_slide(prs, section_title, section_num, accent_color=PRIMARY_COLOR):
    """创建章节过渡页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # 大号序号
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}"
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = "Arial"
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.8), Inches(8), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    
    # 装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.2), Inches(3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    return slide

def add_detail_slide(prs, title, bullet_points, accent_color=PRIMARY_COLOR):
    """创建详细内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_color
    top_bar.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = "Microsoft YaHei"
    
    # 标题下划线
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.03)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = accent_color
    underline.line.fill.background()
    
    # 内容列表
    y_pos = Inches(1.6)
    for point in bullet_points:
        # 圆点
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.6), y_pos + Inches(0.15), Inches(0.12), Inches(0.12)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = accent_color
        bullet.line.fill.background()
        
        # 文本
        content_box = slide.shapes.add_textbox(Inches(0.9), y_pos, Inches(11.5), Inches(0.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.font.name = "Microsoft YaHei"
        
        y_pos += Inches(1.0)
    
    return slide

def create_template_1():
    """模板1：经典紫蓝风格（与原PPT一致）"""
    prs = Presentation()
    set_slide_size(prs)
    
    # 封面
    add_title_slide(prs, "2026年中国高校计算机大赛", "- AIGC创新赛", PRIMARY_COLOR)
    
    # 目录页
    add_content_slide(prs, "目录", ["赛事介绍", "参赛规则", "作品要求", "评审标准"], PRIMARY_COLOR)
    
    # 章节页
    add_section_slide(prs, "赛事介绍", 1, PRIMARY_COLOR)
    
    # 内容页
    add_detail_slide(prs, "赛事背景", [
        "AIGC（人工智能生成内容）技术快速发展",
        "推动高校AI创新应用",
        "培养新一代AI人才"
    ], PRIMARY_COLOR)
    
    return prs

def create_template_2():
    """模板2：科技蓝风格"""
    prs = Presentation()
    set_slide_size(prs)
    
    accent = ACCENT_BLUE
    
    add_title_slide(prs, "科技创新大赛", "Technology Innovation", accent)
    add_content_slide(prs, "目录", ["项目概述", "技术方案", "市场分析", "未来规划"], accent)
    add_section_slide(prs, "项目概述", 1, accent)
    add_detail_slide(prs, "项目背景", [
        "数字化转型浪潮",
        "人工智能技术应用",
        "创新解决方案"
    ], accent)
    
    return prs

def create_template_3():
    """模板3：活力粉风格"""
    prs = Presentation()
    set_slide_size(prs)
    
    accent = ACCENT_PINK
    
    add_title_slide(prs, "创意设计大赛", "Creative Design", accent)
    add_content_slide(prs, "目录", ["设计理念", "视觉系统", "交互体验", "品牌传播"], accent)
    add_section_slide(prs, "设计理念", 1, accent)
    add_detail_slide(prs, "设计背景", [
        "用户体验为核心",
        "视觉美学创新",
        "品牌价值传递"
    ], accent)
    
    return prs

def create_template_4():
    """模板4：生态绿风格"""
    prs = Presentation()
    set_slide_size(prs)
    
    accent = ACCENT_GREEN
    
    add_title_slide(prs, "绿色创新大赛", "Green Innovation", accent)
    add_content_slide(prs, "目录", ["环保理念", "技术方案", "实施计划", "社会效益"], accent)
    add_section_slide(prs, "环保理念", 1, accent)
    add_detail_slide(prs, "项目背景", [
        "可持续发展目标",
        "绿色技术应用",
        "生态环境保护"
    ], accent)
    
    return prs

def create_template_5():
    """模板5：活力橙风格"""
    prs = Presentation()
    set_slide_size(prs)
    
    accent = ACCENT_ORANGE
    
    add_title_slide(prs, "创业计划大赛", "Startup Plan", accent)
    add_content_slide(prs, "目录", ["商业计划", "市场分析", "财务预测", "团队介绍"], accent)
    add_section_slide(prs, "商业计划", 1, accent)
    add_detail_slide(prs, "项目背景", [
        "市场需求分析",
        "商业模式创新",
        "团队优势展示"
    ], accent)
    
    return prs

def main():
    output_dir = r"C:\Users\30742\Desktop\PPT模板"
    os.makedirs(output_dir, exist_ok=True)
    
    templates = [
        ("模板1_经典紫蓝风格.pptx", create_template_1),
        ("模板2_科技蓝风格.pptx", create_template_2),
        ("模板3_活力粉风格.pptx", create_template_3),
        ("模板4_生态绿风格.pptx", create_template_4),
        ("模板5_活力橙风格.pptx", create_template_5),
    ]
    
    for filename, create_func in templates:
        prs = create_func()
        filepath = os.path.join(output_dir, filename)
        prs.save(filepath)
        print(f"已创建: {filepath}")
    
    print(f"\n所有模板已保存到: {output_dir}")

if __name__ == "__main__":
    main()
